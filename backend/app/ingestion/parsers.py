"""MIME-typed document parsing to plain text — dependency-localized (CC-5 #21).

Turns the stored bytes of an uploaded document into plain text for chunking +
embedding. Only the upload allowlist (spec 0004 / #22 ``UPLOAD_ALLOWED_CONTENT_TYPES``)
is supported: PDF, DOCX, PPTX, XLSX, ``text/plain``, ``text/markdown``.

**Dependency localization (ADR-0004 implementation note).** Each format's parser
library (``pypdf`` / ``python-docx`` / ``python-pptx`` / ``openpyxl``) is imported
**lazily, inside its small helper**, so:

* importing this module (or the task module, app boot, most tests) pulls in none
  of the parser libraries — the heavy/optional dependency stays contained here;
* the set of files that touch a given parser is exactly one function, so a
  swap/upgrade is localized.

**Fail-closed (AC-6 / INV-8).** Every failure path is a typed
:class:`DocumentParseError`: an unsupported MIME type, or a corrupt/unreadable
file of a supported type. The Celery task maps that to ``status=failed`` with the
reason — never a silent drop and never an unmapped crash. Parsing is pure CPU
work with no network and no I/O beyond the in-memory bytes it is handed.
"""

from __future__ import annotations

import io

# Upload-allowlist MIME types (kept in lockstep with #22's
# ``Settings.upload_allowed_content_types`` default; the task validates against
# the live setting, this is the parser-dispatch domain).
_PDF = "application/pdf"
_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_TXT = "text/plain"
_MD = "text/markdown"

SUPPORTED_MIME_TYPES: frozenset[str] = frozenset({_PDF, _DOCX, _PPTX, _XLSX, _TXT, _MD})


class DocumentParseError(Exception):
    """Parsing failed — unsupported type or a corrupt/unreadable file (AC-6).

    Carried by the task into ``Document.status=failed`` with the reason. A domain
    error, not an HTTP one; the task records it on the document row.
    """


class UnsupportedMimeTypeError(DocumentParseError):
    """The document's declared MIME type has no parser (outside the allowlist)."""


def _parse_text(data: bytes) -> str:
    """Decode plain text / markdown bytes as UTF-8 (lenient on stray bytes).

    ``errors="replace"`` keeps a stray non-UTF-8 byte from failing the whole
    ingest — the document is still mostly text and worth indexing; the rare
    replacement character is preferable to a hard failure for a text upload.
    """
    return data.decode("utf-8", errors="replace")


# The OOXML formats (docx/pptx/xlsx) are zip containers; corrupt/foreign bytes
# surface as a wide range of library-internal exceptions (BadZipFile, KeyError,
# lxml errors, …). Parsing *untrusted* upload bytes must never crash the worker,
# so each helper catches broadly and re-raises one typed DocumentParseError
# (AC-6); the original is chained for logs. ``MemoryError``/``KeyboardInterrupt``
# (BaseException) are deliberately not caught.


def _parse_pdf(data: bytes) -> str:
    """Extract text from a PDF, page by page (``pypdf``, imported lazily)."""
    from pypdf import PdfReader

    try:
        reader = PdfReader(io.BytesIO(data))
        pages = [page.extract_text() or "" for page in reader.pages]
    except Exception as exc:  # noqa: BLE001 — untrusted bytes; mapped to a typed error
        raise DocumentParseError(f"could not parse PDF: {type(exc).__name__}") from exc
    return "\n\n".join(pages)


def _parse_docx(data: bytes) -> str:
    """Extract paragraph text from a DOCX (``python-docx``, imported lazily)."""
    import docx

    try:
        document = docx.Document(io.BytesIO(data))
        paragraphs = [p.text for p in document.paragraphs]
    except Exception as exc:  # noqa: BLE001 — untrusted bytes; mapped to a typed error
        raise DocumentParseError(f"could not parse DOCX: {type(exc).__name__}") from exc
    return "\n".join(paragraphs)


def _parse_pptx(data: bytes) -> str:
    """Extract text-frame text from a PPTX (``python-pptx``, imported lazily)."""
    from pptx import Presentation

    try:
        presentation = Presentation(io.BytesIO(data))
        lines: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in paragraph.runs)
                        if text:
                            lines.append(text)
    except Exception as exc:  # noqa: BLE001 — untrusted bytes; mapped to a typed error
        raise DocumentParseError(f"could not parse PPTX: {type(exc).__name__}") from exc
    return "\n".join(lines)


def _parse_xlsx(data: bytes) -> str:
    """Extract cell text from an XLSX (``openpyxl``, imported lazily).

    Reads values only (``data_only=True``) and renders each non-empty row as a
    tab-joined line, sheets separated by a blank line. Read-only mode keeps the
    workbook off the heap for large sheets.
    """
    from openpyxl import load_workbook

    try:
        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        sheets: list[str] = []
        for worksheet in workbook.worksheets:
            rows: list[str] = []
            for row in worksheet.iter_rows(values_only=True):
                cells = [str(cell) for cell in row if cell is not None]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                sheets.append("\n".join(rows))
        workbook.close()
    except Exception as exc:  # noqa: BLE001 — untrusted bytes; mapped to a typed error
        raise DocumentParseError(f"could not parse XLSX: {type(exc).__name__}") from exc
    return "\n\n".join(sheets)


# MIME type -> parser. Each value localizes its library import to its own body.
_PARSERS = {
    _PDF: _parse_pdf,
    _DOCX: _parse_docx,
    _PPTX: _parse_pptx,
    _XLSX: _parse_xlsx,
    _TXT: _parse_text,
    _MD: _parse_text,
}


def parse_document(data: bytes, *, mime_type: str) -> str:
    """Extract plain text from ``data`` for the allowlisted ``mime_type``.

    Dispatches on the (normalized) MIME type to the matching helper. The leading
    parameter of a ``Content-Type`` (e.g. ``text/plain; charset=utf-8``) is
    tolerated by splitting on ``;``.

    Raises:
        UnsupportedMimeTypeError: ``mime_type`` is outside the allowlist (AC-6).
        DocumentParseError: a supported type whose bytes are corrupt/unreadable.
    """
    normalized = mime_type.split(";", 1)[0].strip().lower()
    parser = _PARSERS.get(normalized)
    if parser is None:
        raise UnsupportedMimeTypeError(f"unsupported MIME type for ingestion: {normalized!r}")
    return parser(data)
