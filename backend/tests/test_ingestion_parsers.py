"""Parser unit tests — every supported type round-trips, failures are typed.

Tests for :mod:`app.ingestion.parsers` (CC-5 #21, AC-1/AC-6). Each fixture is
generated in-memory with the same OSS library the parser uses, so no binary blobs
are checked in and the suite stays offline. Asserts:

* AC-1: pdf / docx / pptx / xlsx / text/plain / text/markdown each extract their
  text content;
* AC-6 (negative): an unsupported MIME type raises ``UnsupportedMimeTypeError``;
  corrupt bytes of a supported type raise ``DocumentParseError`` (never a crash,
  never a silent empty string masking corruption);
* the content-type parameter (``; charset=...``) and case are tolerated.
"""

from __future__ import annotations

import io

import pytest

from app.ingestion.parsers import (
    SUPPORTED_MIME_TYPES,
    DocumentParseError,
    UnsupportedMimeTypeError,
    parse_document,
)

_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


# --- Fixture builders (in-memory, via the parser's own library) -------------


def _make_pdf(text: str) -> bytes:
    # pypdf can't easily author text PDFs; use reportlab if present, else build a
    # minimal one-page PDF by hand is fragile. We instead author via pypdf's
    # writer with a blank page and rely on a tiny hand-built text stream is hard;
    # simplest robust route: use the `fpdf2`-free path — generate with pypdf is
    # not supported, so we synthesize a minimal valid PDF containing the text.
    # A minimal PDF with a single text-showing content stream:
    content = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode("latin-1")
    objects: list[bytes] = []
    objects.append(b"<< /Type /Catalog /Pages 2 0 R >>")
    objects.append(b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>")
    objects.append(
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>"
    )
    stream_obj = b"<< /Length " + str(len(content)).encode() + b" >>\nstream\n" + content
    objects.append(stream_obj + b"\nendstream")
    objects.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    out = io.BytesIO()
    out.write(b"%PDF-1.4\n")
    offsets: list[int] = []
    for i, obj in enumerate(objects, start=1):
        offsets.append(out.tell())
        out.write(f"{i} 0 obj\n".encode())
        out.write(obj)
        out.write(b"\nendobj\n")
    xref_pos = out.tell()
    out.write(f"xref\n0 {len(objects) + 1}\n".encode())
    out.write(b"0000000000 65535 f \n")
    for off in offsets:
        out.write(f"{off:010d} 00000 n \n".encode())
    out.write(b"trailer\n")
    out.write(f"<< /Size {len(objects) + 1} /Root 1 0 R >>\n".encode())
    out.write(b"startxref\n")
    out.write(f"{xref_pos}\n".encode())
    out.write(b"%%EOF")
    return out.getvalue()


def _make_docx(text: str) -> bytes:
    import docx

    document = docx.Document()
    document.add_paragraph(text)
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _make_pptx(text: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text_frame.text = text
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


def _make_xlsx(value: str) -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    worksheet = workbook.active
    worksheet["A1"] = value
    buf = io.BytesIO()
    workbook.save(buf)
    return buf.getvalue()


# --- Supported types (AC-1) -------------------------------------------------


def test_parse_text_plain() -> None:
    out = parse_document(b"hello plain text", mime_type="text/plain")
    assert out == "hello plain text"


def test_parse_markdown() -> None:
    out = parse_document(b"# Heading\n\nbody", mime_type="text/markdown")
    assert "Heading" in out
    assert "body" in out


def test_parse_text_tolerates_content_type_parameters_and_case() -> None:
    out = parse_document(b"data", mime_type="TEXT/Plain; charset=utf-8")
    assert out == "data"


def test_parse_text_replaces_invalid_utf8_rather_than_crashing() -> None:
    out = parse_document(b"ok \xff\xfe bytes", mime_type="text/plain")
    assert "ok" in out and "bytes" in out  # decoded leniently, no crash


def test_parse_pdf() -> None:
    needle = "Hello PDF world"
    out = parse_document(_make_pdf(needle), mime_type="application/pdf")
    assert needle.replace(" ", "") in out.replace(" ", "")


def test_parse_docx() -> None:
    out = parse_document(_make_docx("Hello DOCX world"), mime_type=_DOCX)
    assert "Hello DOCX world" in out


def test_parse_pptx() -> None:
    out = parse_document(_make_pptx("Hello PPTX world"), mime_type=_PPTX)
    assert "Hello PPTX world" in out


def test_parse_xlsx() -> None:
    out = parse_document(_make_xlsx("Hello XLSX world"), mime_type=_XLSX)
    assert "Hello XLSX world" in out


def test_supported_mime_types_set_is_the_upload_allowlist() -> None:
    assert SUPPORTED_MIME_TYPES == frozenset(
        {
            "application/pdf",
            _DOCX,
            _PPTX,
            _XLSX,
            "text/plain",
            "text/markdown",
        }
    )


# --- Negatives (AC-6) -------------------------------------------------------


def test_unsupported_mime_type_raises() -> None:
    with pytest.raises(UnsupportedMimeTypeError):
        parse_document(b"<svg/>", mime_type="image/svg+xml")


def test_corrupt_pdf_raises_parse_error_not_crash() -> None:
    with pytest.raises(DocumentParseError):
        parse_document(b"%PDF-1.4 not really a pdf", mime_type="application/pdf")


def test_corrupt_docx_raises_parse_error() -> None:
    with pytest.raises(DocumentParseError):
        parse_document(b"PK\x03\x04 not a real docx", mime_type=_DOCX)


def test_corrupt_xlsx_raises_parse_error() -> None:
    with pytest.raises(DocumentParseError):
        parse_document(b"not a workbook", mime_type=_XLSX)


def test_unsupported_is_a_document_parse_error_subclass() -> None:
    # The task catches DocumentParseError; the unsupported case must be caught too.
    assert issubclass(UnsupportedMimeTypeError, DocumentParseError)
